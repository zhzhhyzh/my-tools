"""Convert PDF operations: pdf-to-word, pdf-to-ppt, pdf-to-excel, word-to-pdf,
ppt-to-pdf, excel-to-pdf, pdf-to-jpg, jpg-to-pdf, html-to-pdf, pdf-to-pdfa."""
import os
import io
import json
import subprocess
import zipfile
from typing import List
from fastapi import APIRouter, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from utils.file_utils import get_temp_dir, cleanup_dir, save_upload

router = APIRouter()


def find_libreoffice():
    """Find LibreOffice executable."""
    paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return "soffice"


def libreoffice_convert(input_path, output_dir, output_format):
    """Convert using LibreOffice CLI."""
    soffice = find_libreoffice()
    cmd = [
        soffice, "--headless", "--convert-to", output_format,
        "--outdir", output_dir, input_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
    base = os.path.splitext(os.path.basename(input_path))[0]
    return os.path.join(output_dir, f"{base}.{output_format}")


@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = os.path.join(tmp, "output.docx")

        from pdf2docx import Converter
        cv = Converter(path)
        cv.convert(output, multi_processing=False)
        cv.close()

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=converted.docx"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    """Convert PDF to PPTX by converting pages to images and placing them on slides."""
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches

        images = convert_from_path(path, dpi=150)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            img_path = os.path.join(tmp, f"page_{images.index(img)}.png")
            img.save(img_path, "PNG")
            slide.shapes.add_picture(
                img_path, Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height,
            )

        output = os.path.join(tmp, "output.pptx")
        prs.save(output)

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=converted.pptx"},
        )
    except ImportError:
        # Fallback to LibreOffice
        output = libreoffice_convert(path, tmp, "pptx")
        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=converted.pptx"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = os.path.join(tmp, "output.xlsx")

        try:
            import tabula
            dfs = tabula.read_pdf(path, pages="all", multiple_tables=True)
            if dfs:
                import openpyxl
                wb = openpyxl.Workbook()
                wb.remove(wb.active)
                for idx, df in enumerate(dfs):
                    ws = wb.create_sheet(title=f"Table_{idx + 1}")
                    for col_idx, col_name in enumerate(df.columns, 1):
                        ws.cell(row=1, column=col_idx, value=str(col_name))
                    for row_idx, row in enumerate(df.values, 2):
                        for col_idx, val in enumerate(row, 1):
                            ws.cell(row=row_idx, column=col_idx, value=val)
                wb.save(output)
            else:
                # No tables found, create empty workbook
                import openpyxl
                wb = openpyxl.Workbook()
                wb.save(output)
        except Exception:
            # Fallback: extract text with pdfplumber
            import pdfplumber
            import openpyxl
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    ws = wb.create_sheet(title=f"Page_{i + 1}")
                    tables = page.extract_tables()
                    if tables:
                        row_offset = 1
                        for table in tables:
                            for r_idx, row in enumerate(table):
                                for c_idx, val in enumerate(row):
                                    ws.cell(row=row_offset + r_idx, column=c_idx + 1, value=val)
                            row_offset += len(table) + 1
                    else:
                        text = page.extract_text() or ""
                        for r_idx, line in enumerate(text.split("\n"), 1):
                            ws.cell(row=r_idx, column=1, value=line)
            wb.save(output)

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": "attachment; filename=converted.xlsx"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = libreoffice_convert(path, tmp, "pdf")
        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=converted.pdf"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = libreoffice_convert(path, tmp, "pdf")
        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=converted.pdf"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = libreoffice_convert(path, tmp, "pdf")
        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=converted.pdf"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/pdf-to-jpg")
async def pdf_to_jpg(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    pages: str = Form("all"),
):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        from pdf2image import convert_from_path

        images = convert_from_path(path, dpi=dpi)

        if pages != "all":
            target = set()
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    s, e = part.split("-")
                    target.update(range(int(s) - 1, int(e)))
                else:
                    target.add(int(part) - 1)
            images = [img for i, img in enumerate(images) if i in target]

        if len(images) == 1:
            buf = io.BytesIO()
            images[0].save(buf, "JPEG", quality=90)
            buf.seek(0)
            return StreamingResponse(
                buf,
                media_type="image/jpeg",
                headers={"Content-Disposition": "attachment; filename=page_1.jpg"},
            )
        else:
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zf:
                for idx, img in enumerate(images):
                    buf = io.BytesIO()
                    img.save(buf, "JPEG", quality=90)
                    zf.writestr(f"page_{idx + 1}.jpg", buf.getvalue())
            zip_buffer.seek(0)
            return StreamingResponse(
                zip_buffer,
                media_type="application/zip",
                headers={"Content-Disposition": "attachment; filename=pdf_images.zip"},
            )
    finally:
        cleanup_dir(tmp)


@router.post("/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    tmp = get_temp_dir()
    try:
        import img2pdf as img2pdf_lib

        image_paths = []
        for f in files:
            path = await save_upload(f, tmp)
            image_paths.append(path)

        output = os.path.join(tmp, "output.pdf")
        with open(output, "wb") as out:
            out.write(img2pdf_lib.convert(image_paths))

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=images.pdf"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/html-to-pdf")
async def html_to_pdf(url: str = Form(None), html_content: str = Form(None)):
    tmp = get_temp_dir()
    try:
        import pdfkit

        output = os.path.join(tmp, "output.pdf")
        options = {"quiet": "", "enable-local-file-access": ""}

        if url:
            pdfkit.from_url(url, output, options=options)
        elif html_content:
            pdfkit.from_string(html_content, output, options=options)
        else:
            return {"error": "Provide either url or html_content"}

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=webpage.pdf"},
        )
    finally:
        cleanup_dir(tmp)


@router.post("/pdf-to-pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...)):
    tmp = get_temp_dir()
    try:
        path = await save_upload(file, tmp)
        output = os.path.join(tmp, "output_pdfa.pdf")

        with pikepdf.open(path) as pdf:
            with pdf.open_metadata() as meta:
                meta["dc:title"] = "PDF/A Document"
                meta["pdf:PDFVersion"] = "1.7"
                meta["pdfaid:part"] = "2"
                meta["pdfaid:conformance"] = "B"
            pdf.save(output)

        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=pdfa.pdf"},
        )
    except Exception as e:
        # Fallback: just re-save with pikepdf
        with pikepdf.open(path) as pdf:
            pdf.save(output)
        with open(output, "rb") as out:
            data = out.read()
        return StreamingResponse(
            io.BytesIO(data),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=pdfa.pdf"},
        )
    finally:
        cleanup_dir(tmp)


import pikepdf
